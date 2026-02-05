from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os

# Paths
brain_dir = r"C:\Users\Lenovo\.gemini\antigravity\brain\1c7708c7-5f76-43a6-901c-7650e8c3b1d3"
hero_img = os.path.join(brain_dir, "ppt_hero_image_1769593386803.png")
workflow_img = os.path.join(brain_dir, "workflow_illustration_1769593710716.png")
tech_img = os.path.join(brain_dir, "tech_stack_visual_1769593967015.png")
student_db_img = os.path.join(brain_dir, "student_dashboard_mockup_1769604428790.png")
admin_db_img = os.path.join(brain_dir, "admin_dashboard_mockup_1769604485054.png")
security_img = os.path.join(brain_dir, "security_lock_visual_1769604569728.png")
db_visual_img = os.path.join(brain_dir, "database_er_visual_1769604678617.png")
success_img = os.path.join(brain_dir, "resolution_success_visual_1769604763838.png")

prs = Presentation()

def add_slide(title_text, content_points, image_path=None, layout_index=1):
    slide_layout = prs.slide_layouts[layout_index]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = title_text
    
    if content_points:
        tf = slide.placeholders[1].text_frame
        for point in content_points:
            p = tf.add_paragraph()
            p.text = point
            p.level = 0
            
    if image_path and os.path.exists(image_path):
        # Professional placement: Side-by-side or bottom center usually
        # Here we place it on the right half
        slide.shapes.add_picture(image_path, Inches(6.2), Inches(1.5), width=Inches(3.3))

# Slide 1: Title
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Student Grievance Redressal System"
subtitle.text = "A Secure & Transparent Platform for Campus Voice\n\nFinal Year Project Presentation"
if os.path.exists(hero_img):
    prs.slides[0].shapes.add_picture(hero_img, Inches(0), Inches(0), width=Inches(10))
    # Send to back
    prs.slides[0].shapes[0]._element.getparent().insert(0, prs.slides[0].shapes[-1]._element)

# Slide 2: Introduction
add_slide("Introduction", [
    "A digital bridge between students and administration.",
    "Designed to streamline the process of reporting and resolving grievances.",
    "Focuses on transparency, accountability, and efficiency.",
    "Empowers students to voice concerns without fear (Anonymous options)."
])

# Slide 3: Problem Statement
add_slide("Problem Statement", [
    "Traditional systems are often manual and slow.",
    "Lack of real-time tracking for students.",
    "Potential for miscommunication or lost complaints.",
    "Admin burden in managing paper trails.",
    "No centralized repository for grievance history."
])

# Slide 4: Objectives
add_slide("Objectives", [
    "Automate the grievance redressal lifecycle.",
    "Provide a user-friendly dashboard for both students and admins.",
    "Enable real-time status updates and notifications.",
    "Ensure data security and student privacy.",
    "Generate insights for long-term campus improvement."
])

# Slide 5: Project Architecture
add_slide("Project Architecture", [
    "Frontend: Responsive Glassmorphism UI (HTML/CSS/JS).",
    "Backend: High-performance API using FastAPI (Python).",
    "Database: Lightweight and efficient storage with SQLite.",
    "Security: JWT (JSON Web Tokens) for authentication.",
    "Architecture: RESTful API design."
], workflow_img)

# Slide 6: Technology Stack
add_slide("Technology Stack", [
    "FastAPI: Modern, fast (high-performance) framework for building APIs.",
    "SQLite + SQLAlchemy: Robust database and ORM for data integrity.",
    "Glassmorphism UI: Trendy, professional, and responsive design.",
    "python-jose/passlib: Industry-standard security and hashing."
], tech_img)

# Slide 7: Key Features - Student (Part 1)
add_slide("Student Features - Onboarding", [
    "Secure Registration: Input name, email, and password.",
    "Dashboard: At-a-glance view of active and past grievances.",
    "Profile Management: Maintaining personal details.",
    "Role-based Access: Ensures students only see their own data."
], student_db_img)

# Slide 8: Key Features - Student (Part 2)
add_slide("Student Features - Submission", [
    "Smart Categorization: Selecting the right department (Academic, Hostel, etc.).",
    "Detailed Description: Contextual input for clarity.",
    "Anonymous Mode: Option to report without revealing identity.",
    "Instant Submission: Immediate logging into the system."
])

# Slide 9: Key Features - Student (Part 3)
add_slide("Student Features - Tracking", [
    "Real-time Status: Pending, In Progress, or Resolved.",
    "Admin Remarks: Viewing feedback from the resolution team.",
    "Timeline View: Chronological history of the grievance.",
    "Search & Filter: Easily locate specific complaints."
])

# Slide 10: Key Features - Admin (Part 1)
add_slide("Admin Features - Overview", [
    "Unified Dashboard: Bird's-eye view of all complaints.",
    "Statistics: Monitoring total, pending, and resolved cases.",
    "System Monitoring: Ensuring the platform is healthy.",
    "Secure Login: Dedicated portal for administrators."
], admin_db_img)

# Slide 11: Key Features - Admin (Part 2)
add_slide("Admin Features - Management", [
    "Actionable List: Direct access to student grievances.",
    "Status Control: Updating progress (e.g., 'In Progress').",
    "Remarking System: Providing official responses/decisions.",
    "Accountability: Records of who resolved which complaint."
])

# Slide 12: Database Schema
add_slide("Database Schema", [
    "Users Table: Stores id, name, email, hashed_password, and role.",
    "Complaints Table: Links to users, category, description, status, and remarks.",
    "Relationships: One-to-Many (One student can have multiple complaints).",
    "Data Integrity: Cascading deletes and foreign key constraints."
], db_visual_img)

# Slide 13: Security & Authentication
add_slide("Security & Authentication", [
    "Password Hashing: Using Argon2/Bcrypt for data protection.",
    "JWT Authentication: Secure communication between Frontend and Backend.",
    "CORS Middleware: Controlled access to API endpoints.",
    "Environment Control: Managed configurations through .env."
], security_img)

# Slide 14: Design Philosophy
add_slide("Design Philosophy", [
    "Visual Excellence: Glassmorphism effects for a premium feel.",
    "Responsive Design: Full functionality on mobile, tablet, and desktop.",
    "Accessibility: High contrast and clear typography (Outfit/Roboto).",
    "Simplicity: Guided user flows to reduce cognitive load."
])

# Slide 15: Conclusion & Future Scope
add_slide("Conclusion & Future Scope", [
    "Successfully delivered a modern solution for campus grievances.",
    "Future Scope: AI-driven complaint classification.",
    "Future Scope: Mobile Application (iOS/Android).",
    "Future Scope: Integration with University ERP systems.",
    "Future Scope: Multi-level administrative approval workflow."
], success_img)

output_path = os.path.join(brain_dir, "Student_Grievance_Redressal_System_Enhanced.pptx")
prs.save(output_path)
print(f"PPT saved to {output_path}")
