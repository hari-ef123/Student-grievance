from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os

# Paths
brain_dir = r"C:\Users\Lenovo\.gemini\antigravity\brain\1c7708c7-5f76-43a6-901c-7650e8c3b1d3"

# Images mapping
images = {
    "bg": os.path.join(brain_dir, "ppt_slide_background_template_1769605396334.png"),
    "hero": os.path.join(brain_dir, "ppt_hero_image_1769593386803.png"),
    "intro": os.path.join(brain_dir, "intro_visual_simple_1769605464901.png"),
    "problem": os.path.join(brain_dir, "problem_statement_visual_1769605565466.png"),
    "objectives": os.path.join(brain_dir, "objectives_visual_simple_1769605647420.png"),
    "architecture": os.path.join(brain_dir, "workflow_illustration_1769593710716.png"),
    "tech": os.path.join(brain_dir, "tech_stack_visual_1769593967015.png"),
    "student_db": os.path.join(brain_dir, "student_dashboard_mockup_1769604428790.png"),
    "submission": os.path.join(brain_dir, "student_submission_visual_1769605703317.png"),
    "tracking": os.path.join(brain_dir, "tracking_visual_simple_1769605772743.png"),
    "admin_db": os.path.join(brain_dir, "admin_dashboard_mockup_1769604485054.png"),
    "database": os.path.join(brain_dir, "database_er_visual_1769604678617.png"),
    "security": os.path.join(brain_dir, "security_lock_visual_1769604569728.png"),
    "success": os.path.join(brain_dir, "resolution_success_visual_1769604763838.png")
}

prs = Presentation()

def apply_background(slide):
    if os.path.exists(images["bg"]):
        # Place background at the very back
        pic = slide.shapes.add_picture(images["bg"], 0, 0, width=Inches(10), height=Inches(7.5))
        slide.shapes[0]._element.getparent().insert(0, pic._element)

def add_slide_premium(title_text, content_points, img_key=None):
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    apply_background(slide)
    
    # Title formatting
    title = slide.shapes.title
    title.text = title_text
    
    # Text placeholder formatting (Move to left 60%)
    body_shape = slide.placeholders[1]
    body_shape.left = Inches(0.5)
    body_shape.top = Inches(1.5)
    body_shape.width = Inches(5.5)
    body_shape.height = Inches(5)
    
    tf = body_shape.text_frame
    tf.word_wrap = True
    for point in content_points:
        p = tf.add_paragraph()
        p.text = point
        p.level = 0
        p.font.size = Pt(20)
    
    # Image placement (Right 40%)
    if img_key and img_key in images and os.path.exists(images[img_key]):
        img_path = images[img_key]
        # Align correctly on the right
        slide.shapes.add_picture(img_path, Inches(6.2), Inches(1.8), width=Inches(3.3))

# --- SLIDES ---

# 1. Title Slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
if os.path.exists(images["hero"]):
    prs.slides[0].shapes.add_picture(images["hero"], 0, 0, width=Inches(10), height=Inches(7.5))
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Student Grievance Redressal System"
subtitle.text = "A Secure & Transparent Platform for Campus Voice\nFinal Year Project Presentation"
# Style subtitle
subtitle.font.color.rgb = (255, 255, 255) if os.path.exists(images["hero"]) else (0, 0, 0)

# 2. Introduction
add_slide_premium("Introduction", [
    "Digital bridge between campus and students.",
    "Replaces manual, slow filing processes.",
    "Ensures accountability for every report.",
    "Available 24/7 for student convenience."
], "intro")

# 3. Problem Statement
add_slide_premium("Problem Statement", [
    "Traditional paper systems are often lost.",
    "Lack of real-time tracking for status.",
    "Significant delays in manual processing.",
    "Student hesitation due to privacy risks."
], "problem")

# 4. Objectives
add_slide_premium("System Objectives", [
    "Automate the grievance logic cycle.",
    "Provide instant visibility to students.",
    "Streamline administrative action items.",
    "Maintain strict data privacy and security."
], "objectives")

# 5. Architecture
add_slide_premium("Project Architecture", [
    "RESTful API design with FastAPI.",
    "Decoupled Frontend (HTML/JS/CSS).",
    "Persistent SQLite storage backend.",
    "JWT-based security layer."
], "architecture")

# 6. Tech Stack
add_slide_premium("Technology Stack", [
    "Backend: FastAPI (High Performance).",
    "Database: SQLite with SQLAlchemy ORM.",
    "Frontend: Modern Vanilla JS & CSS3.",
    "Security: Argon2 & JWT Security."
], "tech")

# 7. Student Features - Dashboard
add_slide_premium("Student: Dashboard", [
    "Personalized overview of complaints.",
    "Status breakdown (Resolved/Pending).",
    "Quick access to submit new issues.",
    "Real-time data synchronization."
], "student_db")

# 8. Student Features - Submission
add_slide_premium("Student: Submission", [
    "Simple form with category selection.",
    "Dynamic description fields.",
    "Optional 'Anonymous' reporting mode.",
    "Instant logging with unique IDs."
], "submission")

# 9. Student Features - Tracking
add_slide_premium("Student: Tracking", [
    "Detailed history of all submissions.",
    "Live status updates from admin side.",
    "Viewing admin remarks and feedback.",
    "Searchable grievance archive."
], "tracking")

# 10. Admin Features - Overview
add_slide_premium("Admin: Dashboard", [
    "Global view of institutional issues.",
    "Summary of pending vs resolved tasks.",
    "Critical alert monitoring.",
    "Management of all student reports."
], "admin_db")

# 11. Admin Features - Management (Reusing Admin DB visual)
add_slide_premium("Admin: Action Center", [
    "Process complaints with a single click.",
    "Add official remarks for clarity.",
    "Update grievance state in real-time.",
    "Assign priority to urgent matters."
], "admin_db")

# 12. Database Design
add_slide_premium("Database Schema", [
    "Relational structure (Users & Complaints).",
    "One-To-Many relationship logic.",
    "Encrypted password storage.",
    "Unique foreign key indexing."
], "database")

# 13. Security Protocols
add_slide_premium("Security & Auth", [
    "JWT (JSON Web Token) encryption.",
    "Secure CORS middleware for UI.",
    "Password hashing (Argon2/Bcrypt).",
    "Endpoint protection for Admin routes."
], "security")

# 14. Design Philosophy (Reusing Intro or Hero)
add_slide_premium("Design & UI", [
    "Futuristic Glassmorphism aesthetics.",
    "Responsive design for all devices.",
    "User-centric navigation flow.",
    "Clean, professional typography."
], "hero")

# 15. Conclusion
add_slide_premium("Conclusion", [
    "Succesful digital redressal solution.",
    "100% Transparency in grievance flow.",
    "Scalable architecture for future units.",
    "Ready for campus-wide deployment."
], "success")

output_path = os.path.join(brain_dir, "Premium_Grievance_Presentation.pptx")
prs.save(output_path)
print(f"Premium PPT saved to {output_path}")
